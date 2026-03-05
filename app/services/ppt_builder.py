from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.schemas import PipelineResult


class PPTBuilder:
    @staticmethod
    def build(result: PipelineResult, output_path: Path) -> Path:
        prs = Presentation()

        for slide_item in result.ppt_outline:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_item.title

            body = slide.placeholders[1].text_frame
            body.clear()
            for idx, point in enumerate(slide_item.key_points):
                if idx == 0:
                    body.text = point
                else:
                    p = body.add_paragraph()
                    p.text = point

            layout = slide_item.layout or {}
            layout_type = layout.get("type")
            if layout_type == "table":
                PPTBuilder._draw_table(slide, layout)
            elif layout_type == "process_flow":
                PPTBuilder._draw_process_flow(slide, layout)
            elif layout_type == "org_chart":
                PPTBuilder._draw_org_chart(slide, layout)
            elif layout_type == "timeline":
                PPTBuilder._draw_timeline(slide, layout)
            elif layout_type in {"layers", "layered_blocks"}:
                PPTBuilder._draw_layers(slide, layout)

            notes = slide.notes_slide.notes_text_frame
            notes.text = slide_item.speaker_note

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    @staticmethod
    def _draw_table(slide, layout: dict) -> None:
        columns = layout.get("columns", [])
        rows = layout.get("rows", [])
        if not columns or not rows:
            return

        table_shape = slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=len(columns),
            left=Inches(0.5),
            top=Inches(3.1),
            width=Inches(12.3),
            height=Inches(3.6),
        )
        table = table_shape.table

        for c, col_name in enumerate(columns):
            table.cell(0, c).text = str(col_name)

        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row[: len(columns)]):
                table.cell(r, c).text = str(value)

    @staticmethod
    def _draw_process_flow(slide, layout: dict) -> None:
        steps = layout.get("steps", [])
        if not steps:
            return

        step_w = 1.9
        left = 0.6
        top = 3.7
        for i, step in enumerate(steps[:6]):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + i * (step_w + 0.2)),
                Inches(top),
                Inches(step_w),
                Inches(0.8),
            )
            shape.text = str(step)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(241, 130, 50)
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    Inches(left + i * (step_w + 0.2) + step_w),
                    Inches(top + 0.2),
                    Inches(0.18),
                    Inches(0.4),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(90, 90, 90)

    @staticmethod
    def _draw_org_chart(slide, layout: dict) -> None:
        nodes = layout.get("nodes", [])
        if not nodes:
            return

        root = nodes[0]
        root_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.2), Inches(2.2), Inches(0.8)
        )
        root_shape.text = str(root.get("name", "운영총괄"))
        root_shape.fill.solid()
        root_shape.fill.fore_color.rgb = RGBColor(241, 130, 50)

        children = root.get("children", [])
        for i, child in enumerate(children[:4]):
            x = 1.5 + i * 2.8
            ch = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.6), Inches(2.3), Inches(0.7)
            )
            ch.text = str(child)
            ch.fill.solid()
            ch.fill.fore_color.rgb = RGBColor(64, 64, 64)
            for p in ch.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255, 255, 255)

    @staticmethod
    def _draw_timeline(slide, layout: dict) -> None:
        phases = layout.get("phases", [])
        if not phases:
            return

        for i, phase in enumerate(phases[:4]):
            x = 1.1 + i * 3.0
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.8), Inches(2.6), Inches(1.2)
            )
            box.text = f"{phase.get('name','')}\n{phase.get('detail','')}"
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(30, 144, 255)
            for p in box.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(255, 255, 255)

    @staticmethod
    def _draw_layers(slide, layout: dict) -> None:
        layers = layout.get("layers", [])
        if not layers:
            return

        for i, layer in enumerate(layers[:5]):
            sh = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.2 + i * 0.7), Inches(11.0), Inches(0.55)
            )
            sh.text = str(layer)
            sh.fill.solid()
            color = 220 - i * 20
            sh.fill.fore_color.rgb = RGBColor(color, color, color)
